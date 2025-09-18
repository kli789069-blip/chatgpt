from pathlib import Path
from typing import Iterable, Sequence

from fpdf import FPDF
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document


class UnsupportedFileType(Exception):
    """Raised when a file type cannot be converted to PDF."""


ALLOWED_EXTENSIONS: Sequence[str] = (
    ".docx",
    ".xlsx",
    ".xlsm",
    ".pptx",
    ".txt",
    ".csv",
)


def convert_file_to_pdf(source_path: Path, output_path: Path) -> None:
    extension = source_path.suffix.lower()
    if extension not in ALLOWED_EXTENSIONS:
        raise UnsupportedFileType("暂不支持该文件类型的PDF转换")

    if extension == ".docx":
        _convert_docx_to_pdf(source_path, output_path)
    elif extension in {".xlsx", ".xlsm"}:
        _convert_excel_to_pdf(source_path, output_path)
    elif extension == ".pptx":
        _convert_pptx_to_pdf(source_path, output_path)
    elif extension == ".txt":
        _convert_text_to_pdf(source_path, output_path)
    elif extension == ".csv":
        _convert_csv_to_pdf(source_path, output_path)
    else:
        raise UnsupportedFileType("暂不支持该文件类型的PDF转换")


def _base_pdf(title: str | None = None) -> FPDF:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    if title:
        pdf.set_font("Arial", "B", 16)
        pdf.multi_cell(0, 12, title)
        pdf.ln(4)
        pdf.set_font("Arial", size=12)
    return pdf


def _convert_docx_to_pdf(source_path: Path, output_path: Path) -> None:
    document = Document(source_path)
    pdf = _base_pdf(title=source_path.stem)

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        pdf.multi_cell(0, 8, text or " ")

    pdf.output(str(output_path))


def _convert_excel_to_pdf(source_path: Path, output_path: Path) -> None:
    workbook = load_workbook(source_path, data_only=True, read_only=True)
    pdf = _base_pdf(title=source_path.stem)

    for worksheet in workbook.worksheets:
        pdf.set_font("Arial", "B", 14)
        pdf.multi_cell(0, 10, f"工作表: {worksheet.title}")
        pdf.set_font("Arial", size=11)
        for row in worksheet.iter_rows(values_only=True):
            pdf.multi_cell(0, 7, _format_row(row))
        pdf.ln(5)
        pdf.set_font("Arial", size=12)

    pdf.output(str(output_path))


def _convert_pptx_to_pdf(source_path: Path, output_path: Path) -> None:
    presentation = Presentation(source_path)
    pdf = _base_pdf(title=source_path.stem)

    for index, slide in enumerate(presentation.slides, start=1):
        if index > 1:
            pdf.add_page()
        pdf.set_font("Arial", "B", 14)
        pdf.multi_cell(0, 10, f"幻灯片 {index}")
        pdf.set_font("Arial", size=12)
        lines = _extract_slide_text(slide)
        if not lines:
            pdf.multi_cell(0, 8, "(此页暂无文本内容)")
        for line in lines:
            pdf.multi_cell(0, 8, line)

    pdf.output(str(output_path))


def _convert_text_to_pdf(source_path: Path, output_path: Path) -> None:
    pdf = _base_pdf(title=source_path.stem)
    with source_path.open("r", encoding="utf-8", errors="ignore") as handle:
        for line in handle:
            pdf.multi_cell(0, 8, line.rstrip() or " ")
    pdf.output(str(output_path))


def _convert_csv_to_pdf(source_path: Path, output_path: Path) -> None:
    pdf = _base_pdf(title=source_path.stem)
    with source_path.open("r", encoding="utf-8", errors="ignore") as handle:
        for line in handle:
            pdf.multi_cell(0, 8, line.rstrip())
    pdf.output(str(output_path))


def _format_row(row: Iterable) -> str:
    values = ["" if cell is None else str(cell) for cell in row]
    return " | ".join(values).strip() or " "


def _extract_slide_text(slide) -> list[str]:
    lines: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = shape.text.strip()
        if text:
            lines.extend(part for part in text.splitlines() if part.strip())
    return lines
